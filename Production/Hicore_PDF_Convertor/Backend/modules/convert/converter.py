import os
import logging
import pandas as pd
import pdfplumber
import zipfile
import shutil
from abc import ABC, abstractmethod
from pdf2image import convert_from_path
from pdf2docx import Converter
from PIL import Image
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape, portrait
from reportlab.platypus import SimpleDocTemplate, LongTable, TableStyle, Paragraph
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle


# ============================================================
# 1️⃣ Base Converter Interface
# ============================================================
class ConverterBase(ABC):
    @abstractmethod
    def Convert(self, InputFile: str, OutputFile: str):
        pass


# ============================================================
# 2️⃣ Converter Implementations
# ============================================================

# ------------------------------------------------------------
# PDF → Word
# ------------------------------------------------------------
class PdfToWordConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
        try:
            logging.info(f"Starting PDF → Word conversion for {InputFile}")
            cv = Converter(InputFile)
            cv.convert(OutputFile, start=0, end=None)
            cv.close()
            logging.info(f"✅ PDF → Word saved at {OutputFile}")
            return OutputFile
        except Exception as e:
            logging.error(f"❌ PDF → Word conversion failed: {e}")
            return None


# ------------------------------------------------------------
# Word → PDF  (LibreOffice, cross-platform)
# ------------------------------------------------------------
class WordToPdfConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
            try:
                import docx
                from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
                from reportlab.lib import colors
                from reportlab.lib.pagesizes import A4
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                import subprocess
                import platform
                import shutil
                import tempfile
                from PyPDF2 import PdfMerger

                # Step 1: Convert Word to PDF using LibreOffice
                system = platform.system()
                logging.info(f"🌍 System detected: {system}")
                output_dir = os.path.dirname(OutputFile)
                if system == "Windows":
                    possible_paths = [
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    ]
                    lo_exec = None
                    for p in possible_paths:
                        if os.path.exists(p):
                            lo_exec = p
                            break
                    if not lo_exec:
                        lo_exec = shutil.which("soffice")
                    if not lo_exec:
                        logging.error("❌ LibreOffice not found on Windows.")
                        return None
                else:
                    lo_exec = shutil.which("libreoffice") or shutil.which("soffice")
                    if not lo_exec:
                        logging.error("❌ LibreOffice not installed.")
                        return None
                logging.info(f"🟩 Using LibreOffice binary: {lo_exec}")
                cmd = [
                    lo_exec,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", output_dir,
                    InputFile,
                ]
                result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                if result.returncode != 0:
                    logging.error(f"❌ LibreOffice error: {result.stderr.decode(errors='ignore')}")
                    return None
                base = os.path.splitext(os.path.basename(InputFile))[0]
                generated_pdf = os.path.join(output_dir, f"{base}.pdf")
                if generated_pdf != OutputFile:
                    os.replace(generated_pdf, OutputFile)

                # Step 2: Extract tables and generate tables PDF
                doc = docx.Document(InputFile)
                tables = doc.tables
                if tables:
                    logging.info(f"Tables detected in Word document: {InputFile}")
                    with tempfile.NamedTemporaryFile(delete=False, suffix="_tables.pdf") as tmp_tables_pdf:
                        tables_pdf_path = tmp_tables_pdf.name
                    styles = getSampleStyleSheet()
                    style = ParagraphStyle("tbl", parent=styles["BodyText"], fontSize=8, leading=10)
                    elements = []
                    for idx, table in enumerate(tables, start=1):
                        data = []
                        for row in table.rows:
                            data.append([Paragraph(cell.text, style) for cell in row.cells])
                        tbl = Table(data, repeatRows=1)
                        tbl.setStyle(TableStyle([
                            ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                            ("GRID", (0, 0), (-1, -1), 0.25, colors.black),
                            ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                        ]))
                        elements.append(Paragraph(f"Table {idx}", styles["Heading4"]))
                        elements.append(tbl)
                        elements.append(Spacer(1, 12))
                    pdf = SimpleDocTemplate(tables_pdf_path, pagesize=A4)
                    pdf.build(elements)

                    # Step 3: Merge tables PDF to main PDF
                    merger = PdfMerger()
                    merger.append(OutputFile)
                    merger.append(tables_pdf_path)
                    merger.write(OutputFile)
                    merger.close()
                    os.remove(tables_pdf_path)
                    logging.info(f"✅ Word → PDF (with tables appended) saved at {OutputFile}")
                    return OutputFile

                logging.info(f"✅ Word → PDF saved at {OutputFile}")
                return OutputFile
            except Exception as e:
                logging.error(f"❌ Word → PDF failed: {e}", exc_info=True)
                return None


# ------------------------------------------------------------
# PDF → Excel (stable version)
# ------------------------------------------------------------
class PdfToExcelConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
        try:
            logging.info("📄 Extracting tables using pdfplumber")

            tables = []

            with pdfplumber.open(InputFile) as pdf:
                for page_index, page in enumerate(pdf.pages, start=1):
                    extracted = page.extract_tables()

                    if not extracted:
                        logging.info(f"ℹ️ No tables on page {page_index}")
                        continue

                    for tbl in extracted:
                        cleaned = [
                            row for row in tbl
                            if any(cell not in [None, "", " "] for cell in row)
                        ]
                        if cleaned:
                            tables.append(pd.DataFrame(cleaned))

            # Always generate an Excel file (even if empty)
            if not tables:
                logging.warning("⚠️ No tables found — creating empty Excel file.")
                pd.DataFrame().to_excel(OutputFile, index=False)
                return OutputFile

            with pd.ExcelWriter(OutputFile) as writer:
                for i, df in enumerate(tables, start=1):
                    df.to_excel(writer, sheet_name=f"table_{i}", index=False)

            logging.info(f"✅ PDF → Excel saved at {OutputFile}")
            return OutputFile

        except Exception as e:
            logging.error(f"❌ PDF → Excel conversion failed: {e}", exc_info=True)
            return None


# ------------------------------------------------------------
# Excel → PDF
# ------------------------------------------------------------
class ExcelToPdfConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
        try:
            df = pd.read_excel(InputFile)

            # Auto page size
            page_size = landscape(A4) if len(df.columns) > 5 else portrait(A4)
            pdf = SimpleDocTemplate(OutputFile, pagesize=page_size)
            styles = getSampleStyleSheet()
            style = ParagraphStyle("tbl", parent=styles["BodyText"], fontSize=8, leading=10)

            data = [list(df.columns)] + df.values.tolist()
            wrapped = [[Paragraph(str(cell), style) for cell in row] for row in data]

            table = LongTable(wrapped, repeatRows=1)
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.black),
            ]))

            pdf.build([table])
            logging.info(f"✅ Excel → PDF saved at {OutputFile}")
            return OutputFile

        except Exception as e:
            logging.error(f"❌ Excel → PDF failed: {e}")
            return None


# ------------------------------------------------------------
# PDF → PowerPoint
# ------------------------------------------------------------
class PdfToPowerPointConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
        try:
            images = convert_from_path(InputFile)
            prs = Presentation()

            for i, img in enumerate(images):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                temp = f"ppt_temp_{i+1}.png"
                img.save(temp)
                slide.shapes.add_picture(temp, 0, 0, prs.slide_width, prs.slide_height)
                os.remove(temp)

            prs.save(OutputFile)
            logging.info(f"✅ PDF → PowerPoint saved at {OutputFile}")
            return OutputFile

        except Exception as e:
            logging.error(f"❌ PDF → PowerPoint failed: {e}")
            return None


# ------------------------------------------------------------
# PowerPoint → PDF (LibreOffice)
# ------------------------------------------------------------
class PowerPointToPdfConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
        import subprocess
        import platform
        import shutil

        try:
            system = platform.system()
            output_dir = os.path.dirname(OutputFile)

            if system == "Windows":
                possible_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                ]
                lo_exec = None
                for p in possible_paths:
                    if os.path.exists(p):
                        lo_exec = p
                        break
                if not lo_exec:
                    lo_exec = shutil.which("soffice")
                if not lo_exec:
                    logging.error("❌ LibreOffice not found.")
                    return None
            else:
                lo_exec = shutil.which("libreoffice") or shutil.which("soffice")
                if not lo_exec:
                    logging.error("❌ LibreOffice not installed.")
                    return None

            cmd = [
                lo_exec,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                InputFile,
            ]

            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

            if result.returncode != 0:
                logging.error(f"❌ LibreOffice PPT→PDF error: {result.stderr.decode(errors='ignore')}")
                return None

            base = os.path.splitext(os.path.basename(InputFile))[0]
            generated = os.path.join(output_dir, f"{base}.pdf")

            if generated != OutputFile:
                os.replace(generated, OutputFile)

            logging.info(f"✅ PowerPoint → PDF saved at {OutputFile}")
            return OutputFile

        except Exception as e:
            logging.error(f"❌ PPT → PDF failed: {e}", exc_info=True)
            return None


# ------------------------------------------------------------
# PDF → Image (always return list)
# ------------------------------------------------------------
class PdfToImageConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile):
        try:
            logging.info(f"📄 Converting PDF → images: {InputFile}")

            Image.MAX_IMAGE_PIXELS = None
            pages = convert_from_path(InputFile)

            output_dir = os.path.join(os.path.dirname(OutputFile), "temp_pages")
            os.makedirs(output_dir, exist_ok=True)

            img_paths = []
            base = os.path.splitext(os.path.basename(InputFile))[0]

            for i, p in enumerate(pages, start=1):
                img_path = os.path.join(output_dir, f"{base}_page{i}.jpg")
                p.convert("RGB").save(img_path, "JPEG", quality=85)
                img_paths.append(img_path)

            logging.info(f"✅ PDF → Image created {len(img_paths)} file(s)")
            return img_paths  # ALWAYS list

        except Exception as e:
            logging.error(f"❌ PDF → Image failed: {e}", exc_info=True)
            return None


# ------------------------------------------------------------
# Image → PDF (supports list or single)
# ------------------------------------------------------------
class ImageToPdfConverter(ConverterBase):
    def Convert(self, InputFile, OutputFile, page_format: str = "a4", **kwargs):
        """
        Convert image(s) to PDF
        
        Args:
            InputFile: Image path or list of image paths
            OutputFile: Output PDF path
            page_format: "a4" (default) or "original" (preserve image size)
            **kwargs: Additional parameters (for compatibility)
        
        Returns:
            Path to output PDF file or None on error
        """
        try:
            from PIL import Image as PILImage
            from reportlab.lib.pagesizes import A4, letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            import io
            
            # Normalize to list
            if isinstance(InputFile, str):
                imgs = [InputFile]
            else:
                imgs = InputFile

            # Open all images
            pil_imgs = [PILImage.open(img).convert("RGB") for img in imgs]
            
            if page_format.lower() == "original":
                # Use original image dimensions
                logging.info(f"📐 Using original image dimensions for PDF")
                pil_imgs[0].save(OutputFile, save_all=True, append_images=pil_imgs[1:])
            else:
                # Use A4 format (default)
                logging.info(f"📄 Using A4 format for PDF")
                self._create_a4_pdf(pil_imgs, OutputFile)

            logging.info(f"✅ Image(s) → PDF saved at {OutputFile} (Format: {page_format})")
            return OutputFile

        except Exception as e:
            logging.error(f"❌ Image → PDF failed: {e}", exc_info=True)
            return None
    
    def _create_a4_pdf(self, pil_imgs, output_file):
        """Create PDF with A4 page size, scaling images to fit"""
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.units import inch
        import tempfile
        import os
        
        width, height = A4  # 595 x 842 points
        margin = 0.5 * inch
        
        c = rl_canvas.Canvas(output_file, pagesize=A4)
        
        for idx, pil_img in enumerate(pil_imgs):
            # Get image dimensions
            img_width, img_height = pil_img.size
            aspect_ratio = img_height / img_width
            
            # Calculate scaled dimensions to fit A4 with margins
            max_width = width - (2 * margin)
            max_height = height - (2 * margin)
            
            # Scale to fit
            if aspect_ratio > (max_height / max_width):
                # Height is limiting factor
                scaled_height = max_height
                scaled_width = scaled_height / aspect_ratio
            else:
                # Width is limiting factor
                scaled_width = max_width
                scaled_height = scaled_width * aspect_ratio
            
            # Center on page
            x = (width - scaled_width) / 2
            y = height - margin - scaled_height
            
            # Save PIL image to temp file for reportlab
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                tmp_path = tmp.name
                pil_img.save(tmp_path, format='PNG')
            
            try:
                # Draw image on canvas
                c.drawImage(tmp_path, x, y, width=scaled_width, height=scaled_height)
            finally:
                # Clean up temp file
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
            
            # Add new page if not last image
            if idx < len(pil_imgs) - 1:
                c.showPage()
        
        c.save()
        logging.info(f"✅ Created A4 PDF with {len(pil_imgs)} page(s)")



# ============================================================
# 3️⃣ Converter Factory
# ============================================================
class ConverterFactory:
    @staticmethod
    def GetConverter(ConversionType: str) -> ConverterBase:
        mapping = {
            "pdf_to_word": PdfToWordConverter,
            "word_to_pdf": WordToPdfConverter,
            "pdf_to_excel": PdfToExcelConverter,
            "excel_to_pdf": ExcelToPdfConverter,
            "pdf_to_powerpoint": PdfToPowerPointConverter,
            "powerpoint_to_pdf": PowerPointToPdfConverter,
            "pdf_to_image": PdfToImageConverter,
            "image_to_pdf": ImageToPdfConverter,
        }
        ConverterClass = mapping.get(ConversionType.lower())
        if not ConverterClass:
            raise ValueError(f"Unsupported conversion type: {ConversionType}")
        return ConverterClass()
